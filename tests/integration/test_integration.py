"""End-to-end integration tests for the AIXtract extraction pipeline."""
from __future__ import annotations

import csv
import io
import json
import zipfile
from pathlib import Path

import pytest

from aixtract import (
    AudioConfig,
    ChunkingConfig,
    ContentChunk,
    DocumentMetadata,
    ExtractionConfig,
    ExtractionEngine,
    ExtractionResult,
    OCRConfig,
    OutputFormat,
    __version__,
    extract,
)


# ---------------------------------------------------------------------------
# Helper fixtures
# ---------------------------------------------------------------------------


@pytest.fixture()
def sample_txt(tmp_path: Path) -> Path:
    """Create a sample .txt file."""
    p = tmp_path / "sample.txt"
    p.write_text("Hello, AIXtract! This is a plain text document.")
    return p


@pytest.fixture()
def sample_csv(tmp_path: Path) -> Path:
    """Create a sample .csv file."""
    p = tmp_path / "data.csv"
    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(["Name", "Age", "City"])
    writer.writerow(["Alice", "30", "New York"])
    writer.writerow(["Bob", "25", "London"])
    p.write_text(buf.getvalue())
    return p


@pytest.fixture()
def sample_json(tmp_path: Path) -> Path:
    """Create a sample .json file."""
    p = tmp_path / "data.json"
    data = {"project": "AIXtract", "version": "1.0", "tags": ["extraction", "nlp"]}
    p.write_text(json.dumps(data))
    return p


@pytest.fixture()
def sample_xml(tmp_path: Path) -> Path:
    """Create a sample .xml file."""
    p = tmp_path / "data.xml"
    p.write_text(
        '<?xml version="1.0"?>\n'
        "<catalog>\n"
        "  <book><title>Python Guide</title><author>Jane</author></book>\n"
        "  <book><title>Data Science</title><author>Bob</author></book>\n"
        "</catalog>"
    )
    return p


# ---------------------------------------------------------------------------
# 1. test_extract_txt_file
# ---------------------------------------------------------------------------


class TestExtractTxtFile:
    """Verify the top-level extract() function works with a .txt file."""

    def test_extract_txt_file(self, sample_txt: Path) -> None:
        result = extract(str(sample_txt))

        assert isinstance(result, ExtractionResult)
        assert result.success is True
        assert "AIXtract" in result.content
        assert result.error is None
        assert result.metadata.filename == "sample.txt"
        assert result.metadata.converter_used is not None


# ---------------------------------------------------------------------------
# 2. test_extract_csv_file
# ---------------------------------------------------------------------------


class TestExtractCsvFile:
    """Verify CSV extraction produces a markdown table."""

    def test_extract_csv_file(self, sample_csv: Path) -> None:
        result = extract(str(sample_csv))

        assert result.success is True
        # The CSV converter should generate a markdown table
        assert "|" in result.content_markdown
        assert "Name" in result.content_markdown
        assert "Alice" in result.content_markdown
        assert "Bob" in result.content_markdown
        assert result.metadata.converter_used == "csv"


# ---------------------------------------------------------------------------
# 3. test_extract_json_file
# ---------------------------------------------------------------------------


class TestExtractJsonFile:
    """Verify JSON extraction populates content_json."""

    def test_extract_json_file(self, sample_json: Path) -> None:
        result = extract(str(sample_json))

        assert result.success is True
        assert result.content_json is not None
        assert result.content_json["project"] == "AIXtract"
        assert result.content_json["version"] == "1.0"
        assert result.metadata.converter_used == "json"


# ---------------------------------------------------------------------------
# 4. test_extract_xml_file
# ---------------------------------------------------------------------------


class TestExtractXmlFile:
    """Verify XML extraction extracts text content."""

    def test_extract_xml_file(self, sample_xml: Path) -> None:
        result = extract(str(sample_xml))

        assert result.success is True
        assert "Python Guide" in result.content
        assert "Jane" in result.content
        assert "Data Science" in result.content
        assert result.metadata.converter_used == "xml"


# ---------------------------------------------------------------------------
# 5. test_extract_html_bytes
# ---------------------------------------------------------------------------


class TestExtractHtmlBytes:
    """Verify extraction from raw HTML bytes with a filename hint."""

    def test_extract_html_bytes(self) -> None:
        html_bytes = b"<html><head><title>Test Page</title></head><body><p>Hello HTML</p></body></html>"
        result = extract(html_bytes, filename="test.html")

        assert result.success is True
        assert "Hello HTML" in result.content
        assert result.metadata.filename == "test.html"


# ---------------------------------------------------------------------------
# 6. test_extract_docx_file
# ---------------------------------------------------------------------------


class TestExtractDocxFile:
    """Verify DOCX extraction via engine."""

    def test_extract_docx_file(self, tmp_path: Path) -> None:
        from docx import Document

        doc_path = tmp_path / "document.docx"
        doc = Document()
        doc.add_heading("Integration Test Title", level=1)
        doc.add_paragraph("This paragraph tests DOCX extraction in AIXtract.")
        doc.save(str(doc_path))

        engine = ExtractionEngine()
        result = engine.extract(doc_path)

        assert result.success is True
        assert "Integration Test Title" in result.content
        assert "DOCX extraction" in result.content
        assert result.metadata.converter_used == "docx"


# ---------------------------------------------------------------------------
# 7. test_extract_xlsx_file
# ---------------------------------------------------------------------------


class TestExtractXlsxFile:
    """Verify XLSX extraction via engine."""

    def test_extract_xlsx_file(self, tmp_path: Path) -> None:
        from openpyxl import Workbook

        xlsx_path = tmp_path / "spreadsheet.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.append(["Product", "Price"])
        ws.append(["Widget", "9.99"])
        ws.append(["Gadget", "19.99"])
        wb.save(str(xlsx_path))

        engine = ExtractionEngine()
        result = engine.extract(xlsx_path)

        assert result.success is True
        assert "Product" in result.content
        assert "Widget" in result.content
        assert result.metadata.converter_used == "xlsx"


# ---------------------------------------------------------------------------
# 8. test_extract_pptx_file
# ---------------------------------------------------------------------------


class TestExtractPptxFile:
    """Verify PPTX extraction via engine."""

    def test_extract_pptx_file(self, tmp_path: Path) -> None:
        from pptx import Presentation

        pptx_path = tmp_path / "presentation.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide Title"
        slide.placeholders[1].text = "Slide body content for testing"
        prs.save(str(pptx_path))

        engine = ExtractionEngine()
        result = engine.extract(pptx_path)

        assert result.success is True
        assert "Slide Title" in result.content
        assert "Slide body content" in result.content
        assert result.metadata.converter_used == "pptx"


# ---------------------------------------------------------------------------
# 9. test_extract_zip_file
# ---------------------------------------------------------------------------


class TestExtractZipFile:
    """Verify ZIP archive extraction."""

    def test_extract_zip_file(self, tmp_path: Path) -> None:
        zip_path = tmp_path / "archive.zip"
        with zipfile.ZipFile(str(zip_path), "w") as zf:
            zf.writestr("readme.txt", "This is the readme content.")
            zf.writestr("notes.txt", "Some notes for the project.")

        result = extract(str(zip_path))

        assert result.success is True
        assert "readme content" in result.content
        assert "notes" in result.content.lower()
        assert result.metadata.converter_used == "archive"


# ---------------------------------------------------------------------------
# 10. test_extract_with_chunking
# ---------------------------------------------------------------------------


class TestExtractWithChunking:
    """Verify chunking is applied when ChunkingConfig is enabled."""

    def test_extract_with_chunking(self, tmp_path: Path) -> None:
        long_path = tmp_path / "long_document.txt"
        # Create a document long enough to require multiple chunks.
        # ContentChunker uses token-based sizing (~4 chars/token),
        # so 100 token chunks ~ 400 chars. 200 sentences ~ 5800 chars.
        long_text = "This is sentence number one. " * 200
        long_path.write_text(long_text)

        config = ExtractionConfig(
            chunking=ChunkingConfig(enabled=True, chunk_size=100, overlap=10),
        )
        result = extract(str(long_path), config=config)

        assert result.success is True
        assert len(result.chunks) > 1
        for chunk in result.chunks:
            assert isinstance(chunk, ContentChunk)
            assert chunk.content
            assert chunk.char_start >= 0
            assert chunk.char_end > chunk.char_start


# ---------------------------------------------------------------------------
# 11. test_extract_nonexistent_file
# ---------------------------------------------------------------------------


class TestExtractNonexistentFile:
    """Verify graceful failure for non-existent files."""

    def test_extract_nonexistent_file(self, tmp_path: Path) -> None:
        missing = tmp_path / "does_not_exist.txt"
        result = extract(str(missing))

        assert result.success is False
        assert result.error is not None
        assert "not found" in result.error.lower() or "File not found" in result.error


# ---------------------------------------------------------------------------
# 12. test_extract_oversized_file
# ---------------------------------------------------------------------------


class TestExtractOversizedFile:
    """Verify that oversized files are rejected."""

    def test_extract_oversized_file(self, sample_txt: Path) -> None:
        config = ExtractionConfig(max_file_size_mb=0)
        result = extract(str(sample_txt), config=config)

        assert result.success is False
        assert result.error is not None
        assert "size" in result.error.lower() or "exceeds" in result.error.lower()


# ---------------------------------------------------------------------------
# 13. test_engine_extract_batch
# ---------------------------------------------------------------------------


class TestEngineExtractBatch:
    """Verify batch extraction with ExtractionEngine."""

    def test_engine_extract_batch(self, tmp_path: Path) -> None:
        files = []
        for i in range(3):
            p = tmp_path / f"batch_{i}.txt"
            p.write_text(f"Batch document number {i} content.")
            files.append(p)

        engine = ExtractionEngine()
        results = list(engine.extract_batch(files, show_progress=False))

        assert len(results) == 3
        for source, result in results:
            assert result.success is True
            assert "Batch document" in result.content


# ---------------------------------------------------------------------------
# 14. test_config_for_llm_preset
# ---------------------------------------------------------------------------


class TestConfigForLlmPreset:
    """Verify the for_llm() preset applies chunking during extraction."""

    def test_config_for_llm_preset(self, tmp_path: Path) -> None:
        long_path = tmp_path / "llm_document.txt"
        long_path.write_text("Sentence for LLM testing. " * 300)

        config = ExtractionConfig.for_llm()
        result = extract(str(long_path), config=config)

        assert result.success is True
        # for_llm enables chunking with chunk_size=2000
        assert config.chunking.enabled is True
        assert len(result.chunks) >= 1


# ---------------------------------------------------------------------------
# 15. test_config_for_rag_preset
# ---------------------------------------------------------------------------


class TestConfigForRagPreset:
    """Verify the for_rag() preset works during extraction."""

    def test_config_for_rag_preset(self, tmp_path: Path) -> None:
        rag_path = tmp_path / "rag_document.txt"
        rag_path.write_text("RAG pipeline content sentence. " * 200)

        config = ExtractionConfig.for_rag()
        result = extract(str(rag_path), config=config)

        assert result.success is True
        assert config.output_format == "text"
        assert config.chunking.enabled is True
        assert config.chunking.chunk_size == 512
        assert len(result.chunks) >= 1


# ---------------------------------------------------------------------------
# 16. test_public_api_imports
# ---------------------------------------------------------------------------


class TestPublicApiImports:
    """Verify all public API symbols are importable from aixtract."""

    def test_public_api_imports(self) -> None:
        import aixtract

        expected_symbols = [
            "__version__",
            "extract",
            "ExtractionEngine",
            "ExtractionConfig",
            "ExtractionResult",
            "DocumentMetadata",
            "ContentChunk",
            "OutputFormat",
            "ChunkingConfig",
            "OCRConfig",
            "AudioConfig",
        ]
        for symbol in expected_symbols:
            assert hasattr(aixtract, symbol), f"Missing public API symbol: {symbol}"

        # Verify __all__ matches expected symbols
        for symbol in expected_symbols:
            assert symbol in aixtract.__all__, f"{symbol} not in aixtract.__all__"


# ---------------------------------------------------------------------------
# 17. test_markitdown_fallback
# ---------------------------------------------------------------------------


class TestMarkitdownFallback:
    """Verify MarkItDown fallback when native converter is disabled."""

    def test_markitdown_fallback(self, sample_txt: Path) -> None:
        config = ExtractionConfig(disabled_converters=["txt"])
        result = extract(str(sample_txt), config=config)

        # When native txt converter is disabled, markitdown should handle .txt
        assert result.success is True
        assert result.metadata.converter_used == "markitdown"
        assert "AIXtract" in result.content


# ---------------------------------------------------------------------------
# 18. test_result_to_markdown
# ---------------------------------------------------------------------------


class TestResultToMarkdown:
    """Verify ExtractionResult.to_markdown() returns markdown content."""

    def test_result_to_markdown(self, sample_txt: Path) -> None:
        result = extract(str(sample_txt))

        assert result.success is True
        md = result.to_markdown()
        assert isinstance(md, str)
        assert len(md) > 0
        assert "AIXtract" in md


# ---------------------------------------------------------------------------
# 19. test_result_to_dict
# ---------------------------------------------------------------------------


class TestResultToDict:
    """Verify ExtractionResult.to_dict() returns a dict with expected keys."""

    def test_result_to_dict(self, sample_txt: Path) -> None:
        result = extract(str(sample_txt))

        assert result.success is True
        d = result.to_dict()
        assert isinstance(d, dict)

        # Top-level keys
        expected_keys = {
            "success",
            "content",
            "content_markdown",
            "content_json",
            "metadata",
            "chunks",
            "error",
            "partial_content",
        }
        assert expected_keys.issubset(d.keys())

        # Metadata sub-keys
        meta = d["metadata"]
        assert "filename" in meta
        assert "converter_used" in meta
        assert meta["filename"] == "sample.txt"
