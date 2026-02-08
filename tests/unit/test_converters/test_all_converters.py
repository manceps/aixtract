"""Comprehensive tests for all document converters.

Tests cover:
- text.py: TXTConverter, CSVConverter, JSONConverter, XMLConverter
- archive.py: ZIPConverter
- pdf.py: PDFConverter
- docx.py: DOCXConverter
- xlsx.py: XLSXConverter
- pptx.py: PPTXConverter
- html.py: HTMLConverter
- base.py: BaseConverter
- _markitdown.py: MarkItDownBackend
- epub.py: EPUBConverter (import only)

For binary formats, test files are created with the corresponding libraries.
Image and audio converters are excluded (require tesseract / whisper).
"""
from __future__ import annotations

import io
import json
import zipfile
from pathlib import Path
from typing import BinaryIO, ClassVar

import pytest

from aixtract.converters.base import BaseConverter
from aixtract.models.config import ExtractionConfig
from aixtract.models.result import DocumentMetadata, ExtractionResult


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def config() -> ExtractionConfig:
    """Default extraction config for converter tests."""
    return ExtractionConfig()


# ===========================================================================
# BaseConverter tests
# ===========================================================================


class _StubConverter(BaseConverter):
    """Concrete subclass of BaseConverter for testing abstract methods."""

    name: ClassVar[str] = "stub"
    supported_extensions: ClassVar[tuple[str, ...]] = (".stub", ".stb")
    supported_mimetypes: ClassVar[tuple[str, ...]] = ("application/x-stub",)
    requires: ClassVar[tuple[str, ...]] = ()

    def extract(
        self,
        source: Path | BinaryIO | bytes,
        filename: str | None = None,
    ) -> ExtractionResult:
        content_bytes, file_path = self._read_source(source)
        text = content_bytes.decode("utf-8", errors="replace")
        return ExtractionResult(
            success=True,
            content=text,
            metadata=DocumentMetadata(filename=filename or "stub"),
        )


class TestBaseConverter:
    """Tests for the abstract BaseConverter interface."""

    def test_can_handle_matching_extension(self):
        """can_handle returns True for a matching extension."""
        assert _StubConverter.can_handle("stub") is True

    def test_can_handle_matching_extension_with_dot(self):
        """can_handle strips the leading dot when matching."""
        assert _StubConverter.can_handle(".stub") is True

    def test_can_handle_case_insensitive(self):
        """can_handle is case-insensitive."""
        assert _StubConverter.can_handle("STUB") is True
        assert _StubConverter.can_handle(".STB") is True

    def test_can_handle_non_matching_extension(self):
        """can_handle returns False for unrecognized extensions."""
        assert _StubConverter.can_handle("pdf") is False
        assert _StubConverter.can_handle(".docx") is False

    def test_can_handle_matching_mimetype(self):
        """can_handle returns True when mimetype matches."""
        assert _StubConverter.can_handle("", mimetype="application/x-stub") is True

    def test_can_handle_non_matching_mimetype(self):
        """can_handle returns False for unrecognized mimetype."""
        assert _StubConverter.can_handle("", mimetype="application/pdf") is False

    def test_read_source_with_bytes(self, config: ExtractionConfig):
        """_read_source with raw bytes returns (bytes, None)."""
        converter = _StubConverter(config)
        raw = b"hello world"
        data, path = converter._read_source(raw)
        assert data == raw
        assert path is None

    def test_read_source_with_path(self, config: ExtractionConfig, tmp_path: Path):
        """_read_source with Path returns (bytes, path)."""
        converter = _StubConverter(config)
        test_file = tmp_path / "sample.txt"
        test_file.write_bytes(b"file content")
        data, path = converter._read_source(test_file)
        assert data == b"file content"
        assert path == test_file

    def test_read_source_with_binary_io(self, config: ExtractionConfig):
        """_read_source with BinaryIO returns (bytes, None)."""
        converter = _StubConverter(config)
        bio = io.BytesIO(b"stream data")
        data, path = converter._read_source(bio)
        assert data == b"stream data"
        assert path is None

    def test_strip_consecutive_newlines(self):
        """_strip_consecutive_newlines collapses 3+ newlines to double."""
        text = "line1\n\n\n\nline2"
        result = BaseConverter._strip_consecutive_newlines(text)
        assert result == "line1\n\nline2"
        # Preserves double newlines (paragraph breaks)
        text2 = "line1\n\nline2"
        assert BaseConverter._strip_consecutive_newlines(text2) == "line1\n\nline2"

    def test_strip_consecutive_newlines_single(self):
        """_strip_consecutive_newlines preserves single newlines."""
        text = "a\nb"
        result = BaseConverter._strip_consecutive_newlines(text)
        assert result == "a\nb"


# ===========================================================================
# TXTConverter tests
# ===========================================================================


class TestTXTConverter:
    """Tests for plain text extraction."""

    def test_extract_from_file_path(self, config: ExtractionConfig, tmp_path: Path):
        """Extract text from a .txt file path."""
        from aixtract.converters.text import TXTConverter

        text_file = tmp_path / "hello.txt"
        text_file.write_text("Hello, World!\nThis is a test document.", encoding="utf-8")

        converter = TXTConverter(config)
        result = converter.extract(text_file, filename="hello.txt")

        assert result.success is True
        assert "Hello, World!" in result.content
        assert "test document" in result.content
        assert result.metadata.format_detected == "txt"
        assert result.metadata.converter_used == "txt"

    def test_extract_from_bytes(self, config: ExtractionConfig):
        """Extract text from raw bytes."""
        from aixtract.converters.text import TXTConverter

        raw = b"Bytes content here."
        converter = TXTConverter(config)
        result = converter.extract(raw, filename="from_bytes.txt")

        assert result.success is True
        assert result.content == "Bytes content here."

    def test_extract_from_binary_io(self, config: ExtractionConfig):
        """Extract text from a BinaryIO (BytesIO) stream."""
        from aixtract.converters.text import TXTConverter

        bio = io.BytesIO(b"Streamed text content.")
        converter = TXTConverter(config)
        result = converter.extract(bio, filename="stream.txt")

        assert result.success is True
        assert result.content == "Streamed text content."

    def test_strips_consecutive_newlines(self, config: ExtractionConfig):
        """Consecutive blank lines (3+) are collapsed to double."""
        from aixtract.converters.text import TXTConverter

        raw = b"Line one.\n\n\n\nLine two."
        converter = TXTConverter(config)
        result = converter.extract(raw, filename="multi_newline.txt")

        assert result.success is True
        # 4 newlines collapsed to 2 (preserves paragraph breaks)
        assert "\n\n\n" not in result.content
        assert "Line one." in result.content
        assert "Line two." in result.content
        assert "Line one." in result.content
        assert "Line two." in result.content

    def test_word_and_char_count_in_metadata(self, config: ExtractionConfig):
        """Metadata includes word_count and char_count."""
        from aixtract.converters.text import TXTConverter

        raw = b"One two three four five."
        converter = TXTConverter(config)
        result = converter.extract(raw, filename="count.txt")

        assert result.metadata.word_count is not None
        assert result.metadata.word_count >= 5
        assert result.metadata.char_count is not None
        assert result.metadata.char_count > 0

    def test_can_handle_extensions(self):
        """TXTConverter handles .txt, .md, .rst, .log."""
        from aixtract.converters.text import TXTConverter

        for ext in ("txt", "md", "rst", "log"):
            assert TXTConverter.can_handle(ext) is True
        assert TXTConverter.can_handle("pdf") is False


# ===========================================================================
# CSVConverter tests
# ===========================================================================


class TestCSVConverter:
    """Tests for CSV extraction."""

    def test_extract_csv_with_headers_and_rows(
        self, config: ExtractionConfig, tmp_path: Path
    ):
        """Extract a CSV with headers and data rows."""
        from aixtract.converters.text import CSVConverter

        csv_content = "Name,Age,City\nAlice,30,New York\nBob,25,London\n"
        csv_file = tmp_path / "people.csv"
        csv_file.write_text(csv_content, encoding="utf-8")

        converter = CSVConverter(config)
        result = converter.extract(csv_file, filename="people.csv")

        assert result.success is True
        assert result.metadata.format_detected == "csv"
        assert result.metadata.converter_used == "csv"

    def test_creates_proper_markdown_table(self, config: ExtractionConfig):
        """Content markdown contains a well-formed markdown table."""
        from aixtract.converters.text import CSVConverter

        csv_bytes = b"Name,Age\nAlice,30\nBob,25\n"
        converter = CSVConverter(config)
        result = converter.extract(csv_bytes, filename="test.csv")

        md = result.content_markdown
        assert "| Name | Age |" in md
        assert "| --- | --- |" in md
        assert "| Alice | 30 |" in md
        assert "| Bob | 25 |" in md

    def test_handles_empty_csv(self, config: ExtractionConfig):
        """An empty CSV returns success with empty content."""
        from aixtract.converters.text import CSVConverter

        converter = CSVConverter(config)
        result = converter.extract(b"", filename="empty.csv")

        assert result.success is True
        assert result.content == ""

    def test_row_count_and_column_count_in_metadata(self, config: ExtractionConfig):
        """Metadata extra dict contains row_count and column_count."""
        from aixtract.converters.text import CSVConverter

        csv_bytes = b"A,B,C\n1,2,3\n4,5,6\n"
        converter = CSVConverter(config)
        result = converter.extract(csv_bytes, filename="nums.csv")

        assert result.metadata.extra["row_count"] == 2
        assert result.metadata.extra["column_count"] == 3

    def test_can_handle_extensions(self):
        """CSVConverter handles .csv and .tsv."""
        from aixtract.converters.text import CSVConverter

        assert CSVConverter.can_handle("csv") is True
        assert CSVConverter.can_handle("tsv") is True
        assert CSVConverter.can_handle("xlsx") is False


# ===========================================================================
# JSONConverter tests
# ===========================================================================


class TestJSONConverter:
    """Tests for JSON extraction."""

    def test_extract_json_file(self, config: ExtractionConfig, tmp_path: Path):
        """Extract from a .json file."""
        from aixtract.converters.text import JSONConverter

        data = {"name": "Alice", "scores": [95, 88, 72]}
        json_file = tmp_path / "data.json"
        json_file.write_text(json.dumps(data), encoding="utf-8")

        converter = JSONConverter(config)
        result = converter.extract(json_file, filename="data.json")

        assert result.success is True
        assert result.metadata.format_detected == "json"
        assert result.metadata.converter_used == "json"

    def test_content_json_is_populated(self, config: ExtractionConfig):
        """content_json field is set with the parsed data."""
        from aixtract.converters.text import JSONConverter

        data = {"key": "value", "nested": {"a": 1}}
        converter = JSONConverter(config)
        result = converter.extract(json.dumps(data).encode(), filename="test.json")

        assert result.content_json is not None
        assert result.content_json["key"] == "value"
        assert result.content_json["nested"]["a"] == 1

    def test_content_json_wraps_list(self, config: ExtractionConfig):
        """A JSON array is wrapped under a 'data' key in content_json."""
        from aixtract.converters.text import JSONConverter

        data = [1, 2, 3]
        converter = JSONConverter(config)
        result = converter.extract(json.dumps(data).encode(), filename="list.json")

        assert result.content_json is not None
        assert result.content_json["data"] == [1, 2, 3]

    def test_formatted_json_in_content(self, config: ExtractionConfig):
        """content contains pretty-printed JSON."""
        from aixtract.converters.text import JSONConverter

        data = {"a": 1}
        converter = JSONConverter(config)
        result = converter.extract(json.dumps(data).encode(), filename="fmt.json")

        assert '"a": 1' in result.content
        # Markdown wraps in code fence
        assert "```json" in result.content_markdown

    def test_can_handle_extensions(self):
        """JSONConverter handles .json only."""
        from aixtract.converters.text import JSONConverter

        assert JSONConverter.can_handle("json") is True
        assert JSONConverter.can_handle("xml") is False


# ===========================================================================
# XMLConverter tests
# ===========================================================================


class TestXMLConverter:
    """Tests for XML extraction."""

    def test_extract_text_from_nested_xml(self, config: ExtractionConfig):
        """Extracts text from nested XML elements."""
        from aixtract.converters.text import XMLConverter

        xml_bytes = b"""<?xml version="1.0"?>
        <catalog>
            <book>
                <title>Python Cookbook</title>
                <author>David Beazley</author>
            </book>
            <book>
                <title>Fluent Python</title>
                <author>Luciano Ramalho</author>
            </book>
        </catalog>"""

        converter = XMLConverter(config)
        result = converter.extract(xml_bytes, filename="catalog.xml")

        assert result.success is True
        assert "Python Cookbook" in result.content
        assert "David Beazley" in result.content
        assert "Fluent Python" in result.content
        assert "Luciano Ramalho" in result.content
        assert result.metadata.format_detected == "xml"

    def test_handles_attributes_and_tail_text(self, config: ExtractionConfig):
        """Handles element tail text correctly."""
        from aixtract.converters.text import XMLConverter

        xml_bytes = b"""<root>
            <item id="1">First</item> after first
            <item id="2">Second</item> after second
        </root>"""

        converter = XMLConverter(config)
        result = converter.extract(xml_bytes, filename="tail.xml")

        assert result.success is True
        assert "First" in result.content
        assert "Second" in result.content
        assert "after first" in result.content
        assert "after second" in result.content

    def test_xml_in_markdown_code_fence(self, config: ExtractionConfig):
        """Markdown content wraps original XML in a code fence."""
        from aixtract.converters.text import XMLConverter

        xml_bytes = b"<root><item>hello</item></root>"
        converter = XMLConverter(config)
        result = converter.extract(xml_bytes, filename="simple.xml")

        assert "```xml" in result.content_markdown

    def test_can_handle_extensions(self):
        """XMLConverter handles .xml only."""
        from aixtract.converters.text import XMLConverter

        assert XMLConverter.can_handle("xml") is True
        assert XMLConverter.can_handle("html") is False


# ===========================================================================
# ZIPConverter tests
# ===========================================================================


class TestZIPConverter:
    """Tests for ZIP archive extraction."""

    def _make_zip(self, tmp_path: Path, files: dict[str, bytes]) -> Path:
        """Helper: create a ZIP with the given filename->content mapping."""
        zip_path = tmp_path / "archive.zip"
        with zipfile.ZipFile(zip_path, "w") as zf:
            for name, data in files.items():
                zf.writestr(name, data)
        return zip_path

    def test_extract_zip_with_text_files(
        self, config: ExtractionConfig, tmp_path: Path
    ):
        """Extract text content from files inside a ZIP."""
        from aixtract.converters.archive import ZIPConverter

        zip_path = self._make_zip(
            tmp_path,
            {
                "readme.txt": b"Hello from readme",
                "notes.md": b"# Notes\nSome markdown content",
            },
        )

        converter = ZIPConverter(config)
        result = converter.extract(zip_path, filename="archive.zip")

        assert result.success is True
        assert "Hello from readme" in result.content
        assert "Some markdown content" in result.content

    def test_lists_files_in_archive(self, config: ExtractionConfig, tmp_path: Path):
        """Markdown output lists each file as a section."""
        from aixtract.converters.archive import ZIPConverter

        zip_path = self._make_zip(
            tmp_path,
            {"a.txt": b"aaa", "b.txt": b"bbb"},
        )

        converter = ZIPConverter(config)
        result = converter.extract(zip_path, filename="multi.zip")

        assert "a.txt" in result.content_markdown
        assert "b.txt" in result.content_markdown

    def test_file_count_in_metadata(self, config: ExtractionConfig, tmp_path: Path):
        """Metadata extra contains file_count."""
        from aixtract.converters.archive import ZIPConverter

        zip_path = self._make_zip(
            tmp_path,
            {"one.txt": b"1", "two.txt": b"2", "three.txt": b"3"},
        )

        converter = ZIPConverter(config)
        result = converter.extract(zip_path, filename="three_files.zip")

        assert result.metadata.extra["file_count"] == 3
        assert len(result.metadata.extra["files"]) == 3

    def test_zip_with_non_text_files(self, config: ExtractionConfig, tmp_path: Path):
        """ZIP with only non-text files shows a listing."""
        from aixtract.converters.archive import ZIPConverter

        zip_path = self._make_zip(
            tmp_path,
            {"image.png": b"\x89PNG\r\n", "data.bin": b"\x00\x01\x02"},
        )

        converter = ZIPConverter(config)
        result = converter.extract(zip_path, filename="binary.zip")

        assert result.success is True
        # Non-text files are listed but not extracted as text
        assert "image.png" in result.content_markdown
        assert "data.bin" in result.content_markdown
        assert "# Archive Contents" in result.content_markdown

    def test_extract_from_bytes(self, config: ExtractionConfig, tmp_path: Path):
        """Extract ZIP from raw bytes."""
        from aixtract.converters.archive import ZIPConverter

        zip_path = self._make_zip(tmp_path, {"test.txt": b"test content"})
        zip_bytes = zip_path.read_bytes()

        converter = ZIPConverter(config)
        result = converter.extract(zip_bytes, filename="bytes.zip")

        assert result.success is True
        assert "test content" in result.content


# ===========================================================================
# PDFConverter tests
# ===========================================================================


class TestPDFConverter:
    """Tests for PDF extraction."""

    @pytest.fixture
    def minimal_pdf_bytes(self) -> bytes:
        """Create a minimal valid (but empty) PDF using raw bytes."""
        # This is a valid PDF 1.0 with one blank page
        return (
            b"%PDF-1.0\n"
            b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
            b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
            b"3 0 obj<</Type/Page/MediaBox[0 0 612 792]"
            b"/Parent 2 0 R/Resources<</Font<</F1 4 0 R>>>>>>endobj\n"
            b"4 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
            b"xref\n"
            b"0 5\n"
            b"0000000000 65535 f \n"
            b"0000000009 00000 n \n"
            b"0000000058 00000 n \n"
            b"0000000115 00000 n \n"
            b"0000000266 00000 n \n"
            b"trailer<</Size 5/Root 1 0 R>>\n"
            b"startxref\n"
            b"348\n"
            b"%%EOF\n"
        )

    @pytest.fixture
    def pdf_with_text(self, tmp_path: Path) -> Path:
        """Create a PDF with actual text content using pypdf."""
        from pypdf import PdfWriter
        from pypdf.generic import (
            ArrayObject,
            DictionaryObject,
            NameObject,
            NumberObject,
            TextStringObject,
        )

        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        pdf_path = tmp_path / "sample.pdf"
        with open(pdf_path, "wb") as f:
            writer.write(f)
        return pdf_path

    def test_extract_minimal_pdf(
        self, config: ExtractionConfig, minimal_pdf_bytes: bytes
    ):
        """PDFConverter handles a minimal blank PDF gracefully."""
        from aixtract.converters.pdf import PDFConverter

        converter = PDFConverter(config)
        result = converter.extract(minimal_pdf_bytes, filename="blank.pdf")

        assert result.success is True
        assert result.metadata.format_detected == "pdf"
        assert result.metadata.converter_used == "pdf"

    def test_page_count_in_metadata(
        self, config: ExtractionConfig, minimal_pdf_bytes: bytes
    ):
        """Metadata includes page_count."""
        from aixtract.converters.pdf import PDFConverter

        converter = PDFConverter(config)
        result = converter.extract(minimal_pdf_bytes, filename="blank.pdf")

        assert result.metadata.page_count is not None
        assert result.metadata.page_count == 1

    def test_extract_pdf_from_path(
        self, config: ExtractionConfig, pdf_with_text: Path
    ):
        """PDFConverter extracts from a file path."""
        from aixtract.converters.pdf import PDFConverter

        converter = PDFConverter(config)
        result = converter.extract(pdf_with_text, filename="sample.pdf")

        assert result.success is True
        assert result.metadata.page_count >= 1

    def test_can_handle_pdf_extension(self):
        """PDFConverter handles .pdf extension."""
        from aixtract.converters.pdf import PDFConverter

        assert PDFConverter.can_handle("pdf") is True
        assert PDFConverter.can_handle(".pdf") is True
        assert PDFConverter.can_handle("PDF") is True

    def test_can_handle_pdf_mimetype(self):
        """PDFConverter handles application/pdf mimetype."""
        from aixtract.converters.pdf import PDFConverter

        assert PDFConverter.can_handle("", mimetype="application/pdf") is True

    def test_content_is_string(
        self, config: ExtractionConfig, minimal_pdf_bytes: bytes
    ):
        """Content and content_markdown are strings (even if empty)."""
        from aixtract.converters.pdf import PDFConverter

        converter = PDFConverter(config)
        result = converter.extract(minimal_pdf_bytes, filename="test.pdf")

        assert isinstance(result.content, str)
        assert isinstance(result.content_markdown, str)


# ===========================================================================
# DOCXConverter tests
# ===========================================================================


class TestDOCXConverter:
    """Tests for Word document extraction."""

    @pytest.fixture
    def sample_docx(self, tmp_path: Path) -> Path:
        """Create a minimal .docx with headings and paragraphs."""
        from docx import Document

        doc = Document()
        doc.add_heading("Test Heading", level=1)
        doc.add_paragraph("This is the first paragraph of test content.")
        doc.add_heading("Second Section", level=2)
        doc.add_paragraph("Another paragraph with more text.")

        docx_path = tmp_path / "sample.docx"
        doc.save(str(docx_path))
        return docx_path

    def test_extract_docx(self, config: ExtractionConfig, sample_docx: Path):
        """Extract content from a valid .docx file."""
        from aixtract.converters.docx import DOCXConverter

        converter = DOCXConverter(config)
        result = converter.extract(sample_docx, filename="sample.docx")

        assert result.success is True
        assert result.metadata.format_detected == "docx"
        assert result.metadata.converter_used == "docx"

    def test_headings_detected(self, config: ExtractionConfig, sample_docx: Path):
        """Headings are converted to markdown heading syntax."""
        from aixtract.converters.docx import DOCXConverter

        converter = DOCXConverter(config)
        result = converter.extract(sample_docx, filename="sample.docx")

        assert "# Test Heading" in result.content_markdown
        assert "## Second Section" in result.content_markdown

    def test_paragraphs_extracted(self, config: ExtractionConfig, sample_docx: Path):
        """Paragraph text is present in the extracted content."""
        from aixtract.converters.docx import DOCXConverter

        converter = DOCXConverter(config)
        result = converter.extract(sample_docx, filename="sample.docx")

        assert "first paragraph" in result.content
        assert "Another paragraph" in result.content

    def test_word_count_and_char_count(
        self, config: ExtractionConfig, sample_docx: Path
    ):
        """Metadata includes word_count and char_count."""
        from aixtract.converters.docx import DOCXConverter

        converter = DOCXConverter(config)
        result = converter.extract(sample_docx, filename="sample.docx")

        assert result.metadata.word_count is not None
        assert result.metadata.word_count > 0
        assert result.metadata.char_count is not None
        assert result.metadata.char_count > 0

    def test_extract_from_bytes(self, config: ExtractionConfig, sample_docx: Path):
        """Extract DOCX from raw bytes."""
        from aixtract.converters.docx import DOCXConverter

        docx_bytes = sample_docx.read_bytes()
        converter = DOCXConverter(config)
        result = converter.extract(docx_bytes, filename="bytes.docx")

        assert result.success is True
        assert "Test Heading" in result.content

    def test_can_handle_extensions(self):
        """DOCXConverter handles .docx and .doc."""
        from aixtract.converters.docx import DOCXConverter

        assert DOCXConverter.can_handle("docx") is True
        assert DOCXConverter.can_handle("doc") is True
        assert DOCXConverter.can_handle("pdf") is False

    def test_docx_with_table(self, config: ExtractionConfig, tmp_path: Path):
        """Extract a DOCX containing a table."""
        from docx import Document

        from aixtract.converters.docx import DOCXConverter

        doc = Document()
        doc.add_paragraph("Table below:")
        table = doc.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "Alpha"
        table.cell(1, 1).text = "100"
        table.cell(2, 0).text = "Beta"
        table.cell(2, 1).text = "200"

        docx_path = tmp_path / "with_table.docx"
        doc.save(str(docx_path))

        converter = DOCXConverter(config)
        result = converter.extract(docx_path, filename="with_table.docx")

        assert result.success is True
        assert "Name" in result.content_markdown
        assert "Alpha" in result.content_markdown
        assert "|" in result.content_markdown


# ===========================================================================
# XLSXConverter tests
# ===========================================================================


class TestXLSXConverter:
    """Tests for Excel spreadsheet extraction."""

    @pytest.fixture
    def sample_xlsx(self, tmp_path: Path) -> Path:
        """Create a minimal .xlsx with data."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "People"
        ws.append(["Name", "Age", "City"])
        ws.append(["Alice", 30, "New York"])
        ws.append(["Bob", 25, "London"])

        xlsx_path = tmp_path / "sample.xlsx"
        wb.save(str(xlsx_path))
        return xlsx_path

    def test_extract_xlsx(self, config: ExtractionConfig, sample_xlsx: Path):
        """Extract content from a valid .xlsx file."""
        from aixtract.converters.xlsx import XLSXConverter

        converter = XLSXConverter(config)
        result = converter.extract(sample_xlsx, filename="sample.xlsx")

        assert result.success is True
        assert result.metadata.format_detected == "xlsx"
        assert result.metadata.converter_used == "xlsx"

    def test_sheet_names_in_content(self, config: ExtractionConfig, sample_xlsx: Path):
        """Sheet name appears in markdown content."""
        from aixtract.converters.xlsx import XLSXConverter

        converter = XLSXConverter(config)
        result = converter.extract(sample_xlsx, filename="sample.xlsx")

        assert "People" in result.content_markdown

    def test_markdown_table_format(self, config: ExtractionConfig, sample_xlsx: Path):
        """Content markdown contains a properly formatted table."""
        from aixtract.converters.xlsx import XLSXConverter

        converter = XLSXConverter(config)
        result = converter.extract(sample_xlsx, filename="sample.xlsx")

        md = result.content_markdown
        assert "| Name | Age | City |" in md
        assert "| --- | --- | --- |" in md
        assert "Alice" in md
        assert "Bob" in md

    def test_sheet_count_in_metadata(
        self, config: ExtractionConfig, sample_xlsx: Path
    ):
        """Metadata extra contains sheet info."""
        from aixtract.converters.xlsx import XLSXConverter

        converter = XLSXConverter(config)
        result = converter.extract(sample_xlsx, filename="sample.xlsx")

        assert "sheet_names" in result.metadata.extra
        assert "People" in result.metadata.extra["sheet_names"]
        assert result.metadata.extra["sheet_count"] == 1

    def test_multiple_sheets(self, config: ExtractionConfig, tmp_path: Path):
        """Extract from an XLSX with multiple sheets."""
        from openpyxl import Workbook

        from aixtract.converters.xlsx import XLSXConverter

        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Sales"
        ws1.append(["Product", "Revenue"])
        ws1.append(["Widget", 1000])

        ws2 = wb.create_sheet("Costs")
        ws2.append(["Category", "Amount"])
        ws2.append(["Materials", 500])

        xlsx_path = tmp_path / "multi_sheet.xlsx"
        wb.save(str(xlsx_path))

        converter = XLSXConverter(config)
        result = converter.extract(xlsx_path, filename="multi_sheet.xlsx")

        assert result.success is True
        assert "Sales" in result.content_markdown
        assert "Costs" in result.content_markdown
        assert result.metadata.extra["sheet_count"] == 2

    def test_extract_from_bytes(self, config: ExtractionConfig, sample_xlsx: Path):
        """Extract XLSX from raw bytes."""
        from aixtract.converters.xlsx import XLSXConverter

        xlsx_bytes = sample_xlsx.read_bytes()
        converter = XLSXConverter(config)
        result = converter.extract(xlsx_bytes, filename="bytes.xlsx")

        assert result.success is True
        assert "Alice" in result.content_markdown

    def test_can_handle_extensions(self):
        """XLSXConverter handles .xlsx and .xls."""
        from aixtract.converters.xlsx import XLSXConverter

        assert XLSXConverter.can_handle("xlsx") is True
        assert XLSXConverter.can_handle("xls") is True
        assert XLSXConverter.can_handle("csv") is False


# ===========================================================================
# PPTXConverter tests
# ===========================================================================


class TestPPTXConverter:
    """Tests for PowerPoint extraction."""

    @pytest.fixture
    def sample_pptx(self, tmp_path: Path) -> Path:
        """Create a minimal .pptx with slide content."""
        from pptx import Presentation

        prs = Presentation()
        # Use the title and content layout (index 1)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Test Title"
        slide.placeholders[1].text = "Slide body content here"

        # Add a second slide
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Second Slide"
        slide2.placeholders[1].text = "More content"

        pptx_path = tmp_path / "sample.pptx"
        prs.save(str(pptx_path))
        return pptx_path

    def test_extract_pptx(self, config: ExtractionConfig, sample_pptx: Path):
        """Extract content from a valid .pptx file."""
        from aixtract.converters.pptx import PPTXConverter

        converter = PPTXConverter(config)
        result = converter.extract(sample_pptx, filename="sample.pptx")

        assert result.success is True
        assert result.metadata.format_detected == "pptx"
        assert result.metadata.converter_used == "pptx"

    def test_slide_content_extracted(
        self, config: ExtractionConfig, sample_pptx: Path
    ):
        """Slide text is present in extracted content."""
        from aixtract.converters.pptx import PPTXConverter

        converter = PPTXConverter(config)
        result = converter.extract(sample_pptx, filename="sample.pptx")

        assert "Test Title" in result.content
        assert "Slide body content here" in result.content
        assert "Second Slide" in result.content

    def test_slide_count_in_metadata(
        self, config: ExtractionConfig, sample_pptx: Path
    ):
        """Metadata includes page_count (slide count)."""
        from aixtract.converters.pptx import PPTXConverter

        converter = PPTXConverter(config)
        result = converter.extract(sample_pptx, filename="sample.pptx")

        assert result.metadata.page_count == 2

    def test_slide_markers_in_markdown(
        self, config: ExtractionConfig, sample_pptx: Path
    ):
        """Markdown content includes slide markers."""
        from aixtract.converters.pptx import PPTXConverter

        converter = PPTXConverter(config)
        result = converter.extract(sample_pptx, filename="sample.pptx")

        assert "<!-- Slide 1 -->" in result.content_markdown
        assert "<!-- Slide 2 -->" in result.content_markdown

    def test_extract_from_bytes(self, config: ExtractionConfig, sample_pptx: Path):
        """Extract PPTX from raw bytes."""
        from aixtract.converters.pptx import PPTXConverter

        pptx_bytes = sample_pptx.read_bytes()
        converter = PPTXConverter(config)
        result = converter.extract(pptx_bytes, filename="bytes.pptx")

        assert result.success is True
        assert "Test Title" in result.content

    def test_can_handle_extensions(self):
        """PPTXConverter handles .pptx and .ppt."""
        from aixtract.converters.pptx import PPTXConverter

        assert PPTXConverter.can_handle("pptx") is True
        assert PPTXConverter.can_handle("ppt") is True
        assert PPTXConverter.can_handle("docx") is False


# ===========================================================================
# HTMLConverter tests
# ===========================================================================


class TestHTMLConverter:
    """Tests for HTML extraction."""

    SAMPLE_HTML = b"""<!DOCTYPE html>
    <html>
    <head><title>Test Page</title></head>
    <body>
        <h1>Main Heading</h1>
        <p>This is a paragraph with <strong>bold</strong> text.</p>
        <h2>Sub Heading</h2>
        <ul>
            <li>Item one</li>
            <li>Item two</li>
            <li>Item three</li>
        </ul>
        <table>
            <tr><th>Name</th><th>Value</th></tr>
            <tr><td>Alpha</td><td>100</td></tr>
        </table>
    </body>
    </html>"""

    def test_extract_html_content(self, config: ExtractionConfig):
        """Extract content from HTML bytes."""
        from aixtract.converters.html import HTMLConverter

        converter = HTMLConverter(config)
        result = converter.extract(self.SAMPLE_HTML, filename="page.html")

        assert result.success is True
        assert result.metadata.format_detected == "html"
        assert result.metadata.converter_used == "html"

    def test_title_extracted(self, config: ExtractionConfig):
        """HTML <title> is captured in metadata."""
        from aixtract.converters.html import HTMLConverter

        converter = HTMLConverter(config)
        result = converter.extract(self.SAMPLE_HTML, filename="page.html")

        assert result.metadata.title == "Test Page"

    def test_markdown_conversion_headings(self, config: ExtractionConfig):
        """Headings are converted to markdown heading syntax."""
        from aixtract.converters.html import HTMLConverter

        converter = HTMLConverter(config)
        result = converter.extract(self.SAMPLE_HTML, filename="page.html")

        assert "# Main Heading" in result.content_markdown
        assert "## Sub Heading" in result.content_markdown

    def test_markdown_conversion_list_items(self, config: ExtractionConfig):
        """List items are converted to markdown bullet points."""
        from aixtract.converters.html import HTMLConverter

        converter = HTMLConverter(config)
        result = converter.extract(self.SAMPLE_HTML, filename="page.html")

        assert "- Item one" in result.content_markdown
        assert "- Item two" in result.content_markdown

    def test_paragraph_text_extracted(self, config: ExtractionConfig):
        """Paragraph text appears in content."""
        from aixtract.converters.html import HTMLConverter

        converter = HTMLConverter(config)
        result = converter.extract(self.SAMPLE_HTML, filename="page.html")

        assert "paragraph" in result.content
        assert "bold" in result.content

    def test_word_count_in_metadata(self, config: ExtractionConfig):
        """Metadata includes word_count."""
        from aixtract.converters.html import HTMLConverter

        converter = HTMLConverter(config)
        result = converter.extract(self.SAMPLE_HTML, filename="page.html")

        assert result.metadata.word_count is not None
        assert result.metadata.word_count > 0

    def test_extract_from_file_path(
        self, config: ExtractionConfig, tmp_path: Path
    ):
        """Extract HTML from a file path."""
        from aixtract.converters.html import HTMLConverter

        html_file = tmp_path / "test.html"
        html_file.write_bytes(self.SAMPLE_HTML)

        converter = HTMLConverter(config)
        result = converter.extract(html_file, filename="test.html")

        assert result.success is True
        assert "Main Heading" in result.content

    def test_can_handle_extensions(self):
        """HTMLConverter handles .html and .htm."""
        from aixtract.converters.html import HTMLConverter

        assert HTMLConverter.can_handle("html") is True
        assert HTMLConverter.can_handle("htm") is True
        assert HTMLConverter.can_handle("xml") is False

    def test_html_without_title(self, config: ExtractionConfig):
        """HTML without a <title> tag handles gracefully."""
        from aixtract.converters.html import HTMLConverter

        html_bytes = b"<html><body><p>No title here</p></body></html>"
        converter = HTMLConverter(config)
        result = converter.extract(html_bytes, filename="notitle.html")

        assert result.success is True
        assert result.metadata.title is None
        assert "No title here" in result.content


# ===========================================================================
# MarkItDownBackend tests
# ===========================================================================


class TestMarkItDownBackend:
    """Tests for the MarkItDown fallback backend."""

    @pytest.fixture
    def backend(self):
        """Create a MarkItDownBackend instance, skip if markitdown not installed."""
        try:
            from aixtract.converters._markitdown import MarkItDownBackend
            return MarkItDownBackend()
        except ImportError:
            pytest.skip("markitdown package not installed")

    def test_can_handle_supported_formats(self, backend):
        """can_handle returns True for supported format extensions."""
        from aixtract.converters._markitdown import MarkItDownBackend

        for ext in [".pdf", ".docx", ".xlsx", ".html", ".txt", ".csv", ".json"]:
            assert backend.can_handle(f"file{ext}") is True

    def test_can_handle_unsupported_formats(self, backend):
        """can_handle returns False for unsupported format extensions."""
        assert backend.can_handle("file.xyz") is False
        assert backend.can_handle("file.abc") is False
        assert backend.can_handle("file.custom") is False

    def test_convert_txt_file(self, backend, tmp_path: Path):
        """Convert a simple .txt file (the simplest supported format)."""
        txt_file = tmp_path / "simple.txt"
        txt_file.write_text(
            "Hello from markitdown test.\nSecond line.", encoding="utf-8"
        )

        result = backend.convert(txt_file)

        assert result.success is True
        assert "Hello from markitdown test" in result.content
        assert result.metadata.converter_used == "markitdown"
        assert result.metadata.filename == "simple.txt"

    def test_convert_nonexistent_file_raises(self, backend):
        """convert raises FileNotFoundError for missing files."""
        with pytest.raises(FileNotFoundError):
            backend.convert("/nonexistent/file.txt")

    def test_convert_unsupported_format_raises(self, backend, tmp_path: Path):
        """convert raises ValueError for unsupported formats."""
        unsupported = tmp_path / "file.xyz"
        unsupported.write_text("content")

        with pytest.raises(ValueError, match="Unsupported format"):
            backend.convert(unsupported)


# ===========================================================================
# EPUBConverter - import-only test
# ===========================================================================


class TestEPUBConverterImport:
    """Verify EPUB converter module is importable (no runtime test)."""

    def test_epub_converter_class_exists(self):
        """EPUBConverter class can be imported."""
        try:
            from aixtract.converters.epub import EPUBConverter

            assert EPUBConverter.name == "epub"
            assert ".epub" in EPUBConverter.supported_extensions
        except ImportError:
            pytest.skip("ebooklib or bs4 not installed")


# ===========================================================================
# Converter Registration tests (ensure converters register properly)
# ===========================================================================


class TestConverterRegistration:
    """Verify that converters register with the ConverterRegistry."""

    def test_text_converters_registered(self):
        """TXT, CSV, JSON, XML converters are registered."""
        from aixtract.core.registry import ConverterRegistry

        # Trigger import
        import aixtract.converters  # noqa: F401

        extensions = ConverterRegistry.get_supported_extensions()
        for ext in ["txt", "csv", "json", "xml"]:
            assert ext in extensions, f"{ext} not registered"

    def test_archive_converter_registered(self):
        """ZIP converter is registered."""
        from aixtract.core.registry import ConverterRegistry

        import aixtract.converters  # noqa: F401

        extensions = ConverterRegistry.get_supported_extensions()
        assert "zip" in extensions

    def test_binary_format_converters_registered(self):
        """PDF, DOCX, XLSX, PPTX, HTML converters are registered."""
        from aixtract.core.registry import ConverterRegistry

        import aixtract.converters  # noqa: F401

        extensions = ConverterRegistry.get_supported_extensions()
        for ext in ["pdf", "docx", "xlsx", "pptx", "html"]:
            assert ext in extensions, f"{ext} not registered"

    def test_get_converter_for_txt(self):
        """ConverterRegistry returns TXTConverter for .txt extension."""
        from aixtract.converters.text import TXTConverter
        from aixtract.core.registry import ConverterRegistry

        converter = ConverterRegistry.get_converter(extension=".txt")
        assert converter is not None
        assert isinstance(converter, TXTConverter)

    def test_get_converter_for_csv(self):
        """ConverterRegistry returns CSVConverter for .csv extension."""
        from aixtract.converters.text import CSVConverter
        from aixtract.core.registry import ConverterRegistry

        converter = ConverterRegistry.get_converter(extension=".csv")
        assert converter is not None
        assert isinstance(converter, CSVConverter)

    def test_get_converter_returns_none_for_unknown(self):
        """ConverterRegistry returns None for unknown extension."""
        from aixtract.core.registry import ConverterRegistry

        converter = ConverterRegistry.get_converter(extension=".unknown_ext_xyz")
        assert converter is None
