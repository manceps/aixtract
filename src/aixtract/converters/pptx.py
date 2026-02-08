"""PowerPoint document converter."""
from __future__ import annotations

import io
from pathlib import Path
from typing import BinaryIO, ClassVar

from aixtract.converters.base import BaseConverter
from aixtract.core.registry import ConverterRegistry
from aixtract.models.result import DocumentMetadata, ExtractionResult


@ConverterRegistry.register
class PPTXConverter(BaseConverter):
    """Extract content from PowerPoint presentations."""

    name: ClassVar[str] = "pptx"
    supported_extensions: ClassVar[tuple[str, ...]] = (".pptx", ".ppt")
    supported_mimetypes: ClassVar[tuple[str, ...]] = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    )
    requires: ClassVar[tuple[str, ...]] = ("pptx",)

    def extract(
        self,
        source: Path | BinaryIO | bytes,
        filename: str | None = None,
    ) -> ExtractionResult:
        """Extract content from PowerPoint file."""
        from pptx import Presentation

        content_bytes, file_path = self._read_source(source)
        prs = Presentation(io.BytesIO(content_bytes))

        text_parts = []
        markdown_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            markdown_parts.append(f"<!-- Slide {slide_num} -->")
            slide_texts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)

                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        rows.append([cell.text.strip() for cell in row.cells])
                    if rows:
                        headers = rows[0]
                        markdown_parts.append(
                            "| " + " | ".join(headers) + " |"
                        )
                        markdown_parts.append(
                            "| " + " | ".join(["---"] * len(headers)) + " |"
                        )
                        for row in rows[1:]:
                            markdown_parts.append(
                                "| " + " | ".join(row) + " |"
                            )

            text_parts.extend(slide_texts)
            markdown_parts.extend(slide_texts)
            markdown_parts.append("")

        content = "\n\n".join(text_parts)
        content_markdown = "\n\n".join(markdown_parts)

        metadata = DocumentMetadata(
            filename=filename or (file_path.name if file_path else "presentation.pptx"),
            file_path=file_path,
            format_detected="pptx",
            page_count=len(prs.slides),
            converter_used=self.name,
            word_count=len(content.split()),
            char_count=len(content),
        )

        return ExtractionResult(
            success=True,
            content=content,
            content_markdown=content_markdown,
            metadata=metadata,
        )
